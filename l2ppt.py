#!/usr/bin/python
# Title: Link2PPT
# Description: Automatic PPTX generator used by
# Rochester2600 group to build a list of links
# Each month.
# Usage: ./l2ppt links.csv
# CSV format:
# url, date, authors

from __future__ import absolute_import
from __future__ import division, print_function, unicode_literals
import csv, argparse, subprocess
import logging, sys
import re
import time
import unicodedata

from sumy.parsers.html import HtmlParser
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lsa import LsaSummarizer as Summarizer
from sumy.nlp.stemmers import Stemmer
from sumy.utils import get_stop_words
from HTMLParser import HTMLParser

try:
    import pptx
except:
    print("Running without PPTX support")
    #sys.exit()

try:
    import remark
except:
    print("no remark support loaded")

try:
    import instalink
except:
    print("No Instapaper support")
import urllib2
try:
    from BeautifulSoup import BeautifulSoup
except:
    print("Missing beautifulsoup module")




OUTPUT = ""
CREDS = "./creds"


def main():
    global OUTPUT
    # Handle arguments
    parser = argparse.ArgumentParser(
        description='Convert links to ppt')
    parser.add_argument("-c",
                        dest='csv',
                        help='csv file')
    parser.add_argument('-i',
        dest='icreds',
        help='File with creds for instapaper')
    parser.add_argument('-r',
        dest="remark",
        help="Output to remarkjs markdown. e.g. output.md")
    parser.add_argument('--full',
        help="Download full list from instapaper",
        action="store_true")
    parser.add_argument('-p',
                        dest='ppt',
                        help="Output to powerpoint")

    args = parser.parse_args()

    content = []

    if args.csv:
        content += parse_csv(args.csv)

    if args.icreds:
        if args.full:
            full = True
        else:
            full = False
        content = get_instapaper(args.icreds, full)

    if args.ppt:
        add_slides(content)
    elif args.remark:
        build_remarks(content, args.remark)
    else:
        print("Output type missing. Choose -r or -p")

def build_remarks(content, path):
    r = remark.Remark()
    for slide in content:
        r.add_slide(slide)
    output = r.build()
    ### Convert from unstripped unicode shit
    #re.sub('<[^<]+?>', '', text)
    output = unicodedata.normalize('NFKD', output).encode('ascii', 'ignore')
    cleanoutput = teh_security(output)
    f = open(path, 'w')
    f.writelines(cleanoutput)
    f.close()

def add_slides(lines):
    for line in lines:
        add_slide(line)

def add_slide(line):
    global OUTPUT
    '''Needs a dict of title and url at least
    '''
    ## TODO this should handle instapaper input
    ##  or a CSV. Right now it's only instapaper
    prs = pptx.Presentation(OUTPUT)
    title_slide_layout = prs.slide_layouts[1] # title / content
    slide = prs.slides.add_slide(title_slide_layout)

    # Find the text share to add content
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
    text_frame.clear()
    for para_str in line["highlights"]:
        p = text_frame.add_paragraph()
        p.text = para_str
    text_frame.add_paragraph().text = line["url"]
    text_frame.add_paragraph().text = line["category"]

    # Set title
    title = slide.shapes.title
    title.text = line["title"]

    # Save to output
    prs.save(OUTPUT)


def desperate_summarizer(content):
    """ Sumy implementation that tries to guess
    what the content means """
    LANGUAGE = "english"
    SENTENCES_COUNT = 5
    parser = PlaintextParser.from_string(content, Tokenizer(LANGUAGE))
    stemmer = Stemmer(LANGUAGE)

    summarizer = Summarizer(stemmer)
    summarizer.stop_words = get_stop_words(LANGUAGE)

    highlights = []
    for sentence in summarizer(parser.document, SENTENCES_COUNT):
        highlights.append(sentence._text)
        print(sentence)

    return highlights


def get_instapaper(creds, full=False):
    f = open(creds).read().splitlines()
    ilink = instalink.Instalink(f)
    ilink.login()
    il = ilink.getlinks()
    links = ilink.handlelinks(il)
    # Only get the last 22 days
    if not full:
        days = 22 * 60 * 60 * 24
        content = list(s for s in links if s["time"] > time.time() - 2595600)
    else:
        content = list(s for s in links)

    for indx, line in enumerate(content):
        if not line["highlights"]:
            logging.debug("No highlights found. adding some")
            print("line[bookmarkid]= %s" % line["bookmark_id"])
            text = ilink.gettext(line["bookmark_id"])
            #text = unicodedata.normalize('NFKD', text).encode('ascii', 'ignore')
            content[indx]["highlights"] = desperate_summarizer(text)


    return content

def teh_security(badness):
    s = Stripper()
    s.feed(badness)
    goodness = s.get_data()
    return goodness


def get_title(url):
    try:
        # 3s timeout
        f = urllib2.urlopen(url, tmeout=3000)
        soup = BeautifulSoup(f)
        f.close()

        if soup.title.string:
            logging.debug("Title found as: %s" % soup.title.string)
            return soup.title.string
        else:
            return "No title found"
    except:
        logging.error("URL: %s had an error" % url)
        return "Blank"


def parse_csv(file):
    # is csv file?
    #csvobj = []
    with open(file, 'rb') as csvfile:
        urllist = csv.reader(csvfile, delimiter=',', quotechar='|')
        content = []
        for row in urllist:
            ## If url TODO
            record = {}
            record["url"] = row[0]
            try:
                record["author"] = row[1]
            except IndexError:
                record["author"] = None

            try:
                record["date"] = row[2]
            except IndexError:
                record["date"] = None
            #sspath = screenshot(record["url"]) ## Disable mem issue
            sspath = None

            if sspath:
                logging.debug("Screenshot complete")
                record["screenshot"] = sspath
                #row.append(sspath)
            else:
                logging.error("Screenshot failed")
                record["screenshot"] = None
            record["title"] = get_title(row[0])

            ## add to slide
            content.append(record)
            #add_slide(record)
        return content



def screenshot(url):
    # screenshot url
    ## http://wkhtmltopdf.org/
    outputfile = "%s.pdf" % url  # TODO strip!!
    commands = ["wkhtmltopdf", url, outputfile]
    try:
        if subprocess.call(commands) < 1:
            logging.debug(
                "Error executing shell command: \n"
                "Make sure you are running with an X session"
            )
            return None
        else:
            return outputfile
    except:
        logging.debug("Exception caught for wkhtmltopdf")
        return None


class Stripper(HTMLParser):
    '''clASS TO summon thE STRpper for tEH STR1PIN'''
    def __init__(self):
        self.reset()
        self.fed = []
    def handle_data(self, d):
        self.fed.append(d)
    def get_data(self):
        return ''.join(self.fed)

if __name__ == '__main__':
    # init main
    main()
